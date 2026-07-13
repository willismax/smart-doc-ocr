"""Layer 3 — 辨識層：各格式的萃取 handler + OCR 引擎封裝。

穩定性設計：
- PaddleOCR lazy 載入且為程序內單例（載入一次 3~10 秒，不能每份文件重載）。
- 掃描 PDF 用 PyMuPDF 內建 rasterizer 轉圖（page.get_pixmap），
  不用 pdf2image → 免裝 poppler，Windows 上少掉最常見的環境問題。
- 大 PDF 分批開頁 + 每批 gc，避免 OOM。
- Office 讀取三段降級：MarkItDown → python-docx/openpyxl → zipfile 抽 XML 文字。
"""
from __future__ import annotations

import gc
import logging
import re
import zipfile
from pathlib import Path

from .config import SETTINGS
from .errors import (CorruptedFileError, OcrEngineUnavailableError,
                     PasswordProtectedError, UnsupportedFormatError)
from .preprocessor import ImagePreprocessor
from .router import DocType

logger = logging.getLogger(__name__)


# ── OCR 引擎（單例、lazy） ─────────────────────────────────────

class OcrEngine:
    """PaddleOCR 封裝。available 為 False 時，數位 PDF / Office 流程不受影響。"""

    _instance = None

    def __init__(self, settings: dict | None = None):
        self.cfg = (settings or SETTINGS)["ocr"]
        self._ocr = None
        self._load_error: str | None = None
        self._loaded = False

    @classmethod
    def get(cls) -> "OcrEngine":
        if cls._instance is None:
            cls._instance = cls()
        return cls._instance

    @property
    def available(self) -> bool:
        self._ensure_loaded()
        return self._ocr is not None

    @property
    def load_error(self) -> str | None:
        self._ensure_loaded()
        return self._load_error

    def _ensure_loaded(self):
        if self._loaded:
            return
        self._loaded = True
        if self.cfg.get("engine") == "none":
            self._load_error = "設定檔停用 OCR 引擎"
            return
        # Windows DLL 相容性：torch 必須先於 paddle 載入，否則之後任何
        # 功能（語意比對）載入 torch 會因 DLL 衝突崩潰（WinError 127）
        try:
            import torch  # noqa: F401
        except Exception:
            pass  # 沒裝 torch（或載入失敗）不影響 OCR 本身
        try:
            from paddleocr import PaddleOCR
        except ImportError as e:
            self._load_error = f"PaddleOCR 未安裝：{e}"
            logger.warning(self._load_error)
            return
        try:
            use_gpu = self.cfg.get("use_gpu", "auto")
            if use_gpu == "auto":
                use_gpu = self._gpu_usable()
            kwargs = dict(
                use_angle_cls=True,
                lang=self.cfg.get("lang", "chinese_cht"),
                show_log=False,
            )
            # 有本地模型目錄就用（完全離線）；沒有讓 PaddleOCR 用預設快取。
            # Paddle 的 C++ 引擎打不開含中文等非 ASCII 字元的路徑（Windows
            # 已知限制），必要時把模型鏡像到 %LOCALAPPDATA% 再載入。
            models = self._ascii_safe_models(Path(SETTINGS["paths"]["models"]))
            model_kwargs = {}
            if models is not None:
                for key, sub in (("det_model_dir", "det"),
                                 ("rec_model_dir", "rec"),
                                 ("cls_model_dir", "cls")):
                    p = models / sub
                    if (p / "inference.pdmodel").is_file():
                        model_kwargs[key] = str(p)
            try:
                self._ocr = self._build_paddle(
                    PaddleOCR, {**kwargs, **model_kwargs}, use_gpu)
            except Exception:
                if not model_kwargs:
                    raise
                # 指定模型目錄載入失敗 → 退回 PaddleOCR 預設快取再試一次
                logger.warning("以 models/ 目錄載入失敗，改用預設快取重試")
                self._ocr = self._build_paddle(PaddleOCR, kwargs, use_gpu)
            logger.info("PaddleOCR 載入完成（GPU=%s）", use_gpu)
        except Exception as e:
            self._load_error = f"PaddleOCR 載入失敗：{e}"
            logger.exception(self._load_error)

    @staticmethod
    def _build_paddle(PaddleOCR, kwargs: dict, use_gpu):
        try:
            return PaddleOCR(use_gpu=bool(use_gpu), **kwargs)
        except TypeError:
            # PaddleOCR 3.x 移除 use_gpu / show_log 參數
            kwargs = {k: v for k, v in kwargs.items() if k != "show_log"}
            return PaddleOCR(**kwargs)

    @staticmethod
    def _ascii_safe_models(models: Path) -> Path | None:
        """回傳 Paddle 可以開啟的模型目錄；來源不存在回傳 None。"""
        if not models.is_dir():
            return None
        if all(ord(c) < 128 for c in str(models)):
            return models
        import os
        import shutil
        base = os.environ.get("LOCALAPPDATA") or str(Path.home())
        mirror = Path(base) / "smart-doc-ocr" / "models"
        if any(ord(c) >= 128 for c in str(mirror)):
            logger.warning("使用者目錄也含非 ASCII 字元，無法鏡像模型：%s", mirror)
            return None
        try:
            for sub in ("det", "rec", "cls"):
                src = models / sub
                dst = mirror / sub
                if (src / "inference.pdmodel").is_file() and \
                        not (dst / "inference.pdmodel").is_file():
                    dst.parent.mkdir(parents=True, exist_ok=True)
                    if dst.exists():
                        shutil.rmtree(dst)
                    shutil.copytree(src, dst)
            logger.info("模型路徑含非 ASCII 字元，已鏡像到 %s", mirror)
            return mirror
        except OSError as e:
            logger.warning("模型鏡像失敗（%s），改用 PaddleOCR 預設快取", e)
            return None

    @staticmethod
    def _gpu_usable() -> bool:
        try:
            import paddle
            return bool(paddle.device.is_compiled_with_cuda())
        except Exception:
            return False

    def ocr_image(self, img) -> str:
        """對 numpy 影像做 OCR，回傳按閱讀順序串接的文字。"""
        self._ensure_loaded()
        if self._ocr is None:
            raise OcrEngineUnavailableError(detail=self._load_error or "")
        result = self._ocr.ocr(img)
        lines: list[str] = []
        if result:
            for page in result:
                if not page:
                    continue
                # PaddleOCR 2.x: [[box, (text, conf)], ...]
                # PaddleOCR 3.x: dict with rec_texts
                if isinstance(page, dict):
                    lines.extend(page.get("rec_texts", []))
                else:
                    for item in page:
                        try:
                            lines.append(item[1][0])
                        except (IndexError, TypeError):
                            continue
        return "\n".join(lines)


# ── 各格式 handler ─────────────────────────────────────────────

class DocumentRecognizer:
    def __init__(self, settings: dict | None = None):
        self.settings = settings or SETTINGS
        self.preprocessor = ImagePreprocessor()
        self.ocr = OcrEngine.get()

    # 主分派
    def extract(self, path: Path, doc_type: DocType) -> dict:
        path = Path(path)
        handlers = {
            DocType.PDF_DIGITAL: self.extract_pdf_digital,
            DocType.PDF_SCANNED: self.extract_pdf_scanned,
            DocType.PDF_MIXED: self.extract_pdf_mixed,
            DocType.WORD: self.extract_office,
            DocType.EXCEL: self.extract_office,
            DocType.PPT: self.extract_office,
            DocType.IMAGE: self.extract_image,
            DocType.EMAIL: self.extract_email,
            DocType.TEXT: self.extract_text,
            DocType.LEGACY_OFFICE: self.extract_office,
        }
        handler = handlers.get(doc_type)
        if handler is None:
            raise UnsupportedFormatError(
                detail=f"{path.name}（{doc_type.value}）")
        return handler(path)

    # ── PDF ──────────────────────────────────────────────────
    def extract_pdf_digital(self, path: Path) -> dict:
        import fitz
        doc = self._open_pdf(path)
        try:
            chunk = self.settings["limits"]["pdf_chunk_pages"]
            total = min(len(doc), self.settings["limits"]["max_pdf_pages"])
            pages = []
            for start in range(0, total, chunk):
                for i in range(start, min(start + chunk, total)):
                    page = doc[i]
                    tables = []
                    try:
                        for t in page.find_tables():
                            tables.append(t.extract())
                    except Exception:
                        pass  # 表格偵測失敗不影響文字萃取
                    pages.append({
                        "page": i + 1,
                        "text": page.get_text().strip(),
                        "tables": tables,
                    })
                gc.collect()
            return {"type": "pdf_digital", "pages": pages, "path": str(path)}
        finally:
            doc.close()

    def extract_pdf_scanned(self, path: Path) -> dict:
        return self._extract_pdf_ocr(path, ocr_all=True)

    def extract_pdf_mixed(self, path: Path) -> dict:
        """逐頁判斷：有文字層直接取，沒有才 OCR。"""
        return self._extract_pdf_ocr(path, ocr_all=False)

    def _extract_pdf_ocr(self, path: Path, ocr_all: bool) -> dict:
        import fitz
        import numpy as np
        doc = self._open_pdf(path)
        min_text = self.settings["ocr"]["min_text_len_per_page"]
        dpi = self.settings["ocr"]["dpi"]
        try:
            total = min(len(doc), self.settings["limits"]["max_pdf_pages"])
            pages = []
            for i in range(total):
                page = doc[i]
                text = page.get_text().strip()
                if not ocr_all and len(text) >= min_text:
                    pages.append({"page": i + 1, "text": text,
                                  "tables": [], "ocr": False})
                    continue
                # rasterize（PyMuPDF 內建，免 poppler）
                pix = page.get_pixmap(dpi=dpi, colorspace=fitz.csGRAY)
                img = np.frombuffer(pix.samples, dtype=np.uint8)
                img = img.reshape(pix.height, pix.width)
                del pix
                processed = self.preprocessor.process(img)
                processed = self.preprocessor.scale_for_ocr(processed)
                ocr_text = self.ocr.ocr_image(processed)
                del img, processed
                pages.append({"page": i + 1, "text": ocr_text,
                              "tables": [], "ocr": True})
                if (i + 1) % self.settings["limits"]["pdf_chunk_pages"] == 0:
                    gc.collect()
            kind = "pdf_scanned" if ocr_all else "pdf_mixed"
            return {"type": kind, "pages": pages, "path": str(path)}
        finally:
            doc.close()
            gc.collect()

    def _open_pdf(self, path: Path):
        import fitz
        try:
            doc = fitz.open(str(path))
        except Exception as e:
            raise CorruptedFileError(detail=f"{path.name}: {e}") from e
        if doc.needs_pass:
            doc.close()
            raise PasswordProtectedError(detail=path.name)
        return doc

    # ── Office ───────────────────────────────────────────────
    def extract_office(self, path: Path) -> dict:
        text = None
        # 第一段：MarkItDown（結構最完整）
        try:
            from markitdown import MarkItDown
            result = MarkItDown(enable_plugins=False).convert(str(path))
            text = result.text_content
        except ImportError:
            logger.info("MarkItDown 未安裝，改用降級讀取器")
        except Exception as e:
            logger.warning("MarkItDown 轉換失敗（%s），改用降級讀取器", e)
        # 第二段：格式專屬套件
        if not text:
            text = self._office_fallback(path)
        # 第三段：直接抽 OOXML 內的 XML 文字
        if not text:
            text = self._ooxml_raw_text(path)
        if text is None:
            raise CorruptedFileError(
                detail=f"{path.name}: 所有 Office 讀取器都失敗")
        return {"type": "office", "text": text, "path": str(path)}

    @staticmethod
    def _office_fallback(path: Path) -> str | None:
        suffix = path.suffix.lower()
        try:
            if suffix == ".docx":
                import docx
                d = docx.Document(str(path))
                parts = [p.text for p in d.paragraphs]
                for table in d.tables:
                    for row in table.rows:
                        parts.append(" | ".join(c.text for c in row.cells))
                return "\n".join(parts)
            if suffix in (".xlsx", ".xlsm"):
                import openpyxl
                wb = openpyxl.load_workbook(str(path), read_only=True,
                                            data_only=True)
                parts = []
                for ws in wb.worksheets:
                    parts.append(f"## 工作表：{ws.title}")
                    for row in ws.iter_rows(values_only=True):
                        cells = ["" if c is None else str(c) for c in row]
                        if any(cells):
                            parts.append(" | ".join(cells))
                wb.close()
                return "\n".join(parts)
            if suffix == ".pptx":
                from pptx import Presentation
                prs = Presentation(str(path))
                parts = []
                for n, slide in enumerate(prs.slides, 1):
                    parts.append(f"## 投影片 {n}")
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            parts.append(shape.text_frame.text)
                return "\n".join(parts)
        except ImportError:
            return None
        except Exception as e:
            logger.warning("降級讀取 %s 失敗：%s", path.name, e)
            return None
        return None

    @staticmethod
    def _ooxml_raw_text(path: Path) -> str | None:
        """最後手段：直接從 OOXML zip 抽出所有 XML 的文字節點。"""
        try:
            with zipfile.ZipFile(path) as zf:
                texts = []
                for name in zf.namelist():
                    if not name.endswith(".xml"):
                        continue
                    if not re.search(r"(document|sheet\d*|slide\d+|sharedStrings)",
                                     name):
                        continue
                    xml = zf.read(name).decode("utf-8", errors="ignore")
                    chunks = re.findall(r">([^<>]+)<", xml)
                    texts.extend(c.strip() for c in chunks if c.strip())
                return "\n".join(texts) if texts else None
        except Exception:
            return None

    # ── 圖片 ─────────────────────────────────────────────────
    def extract_image(self, path: Path) -> dict:
        import numpy as np
        from PIL import Image
        try:
            with Image.open(path) as pil:
                pil = pil.convert("L")  # 灰階
                img = np.array(pil)
        except Exception as e:
            raise CorruptedFileError(detail=f"{path.name}: {e}") from e
        processed = self.preprocessor.process(img)
        processed = self.preprocessor.scale_for_ocr(processed)
        text = self.ocr.ocr_image(processed)
        return {"type": "image", "text": text, "path": str(path)}

    # ── Email ────────────────────────────────────────────────
    def extract_email(self, path: Path) -> dict:
        if path.suffix.lower() == ".msg":
            return self._extract_msg(path)
        import email
        import email.policy
        with open(path, "rb") as f:
            msg = email.message_from_binary_file(f, policy=email.policy.default)
        parts = [f"寄件者：{msg.get('From', '')}",
                 f"收件者：{msg.get('To', '')}",
                 f"主旨：{msg.get('Subject', '')}",
                 f"日期：{msg.get('Date', '')}", ""]
        body = msg.get_body(preferencelist=("plain", "html"))
        if body is not None:
            content = body.get_content()
            if body.get_content_type() == "text/html":
                content = re.sub(r"<[^>]+>", " ", content)
            parts.append(content)
        attachments = [att.get_filename() or "(未命名)"
                       for att in msg.iter_attachments()]
        if attachments:
            parts.append("\n附件：" + "、".join(attachments))
        return {"type": "email", "text": "\n".join(parts), "path": str(path)}

    @staticmethod
    def _extract_msg(path: Path) -> dict:
        try:
            import extract_msg
        except ImportError as e:
            raise UnsupportedFormatError(
                "讀取 Outlook .msg 需要安裝 extract-msg 套件。",
                detail=str(e)) from e
        msg = extract_msg.Message(str(path))
        try:
            text = (f"寄件者：{msg.sender}\n收件者：{msg.to}\n"
                    f"主旨：{msg.subject}\n日期：{msg.date}\n\n{msg.body or ''}")
        finally:
            msg.close()
        return {"type": "email", "text": text, "path": str(path)}

    # ── 純文字 ───────────────────────────────────────────────
    @staticmethod
    def extract_text(path: Path) -> dict:
        raw = path.read_bytes()
        for enc in ("utf-8-sig", "utf-8", "big5", "cp950"):
            try:
                return {"type": "text", "text": raw.decode(enc),
                        "path": str(path)}
            except (UnicodeDecodeError, LookupError):
                continue
        return {"type": "text",
                "text": raw.decode("utf-8", errors="replace"),
                "path": str(path)}
